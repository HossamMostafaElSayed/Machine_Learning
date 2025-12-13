import io
from typing import Optional, Tuple

import cv2
import numpy as np
import yt_dlp
import streamlit as st
from pptx import Presentation
from pptx.util import Inches

def extract_slides_from_url(
    video_url: str,
    skip_seconds: float = 2.0,
    mse_threshold: float = 2000.0,
    max_frames: int = 15000,
) -> Tuple[Optional[io.BytesIO], Optional[str]]:
    """
    Stream a YouTube video (no full download) and extract distinct slide images.

    Args:
        video_url: YouTube video URL or short link.
        skip_seconds: Seconds to skip between sampled frames.
        mse_threshold: Mean-squared-error threshold to consider a frame a new slide.
        max_frames: Safety limit on number of frames processed.

    Returns:
        Tuple of (BytesIO with PPTX data or None, error message or None).
    """

    # 1. Resolve a direct stream URL (yt_dlp will not download media when download=False)
    ydl_opts = {
        "format": "bestvideo+bestaudio/best",
        "noplaylist": True,
        "quiet": True,
        "no_warnings": True,
    }

    try:
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(video_url, download=False)
            # yt_dlp may return a dict or a list for formats; prefer the 'url' top-level if present
            stream_url = info.get("url")
            duration = info.get("duration")
            formats = info.get("formats") or []
    except Exception as e:
        return None, f"Error getting video stream: {str(e)}"

    # If no top-level url was provided, pick a suitable format from `formats`
    if not stream_url:
        # Prefer mp4/http formats with both audio and video if available
        candidates = []
        for f in formats:
            if not f.get("url"):
                continue
            # prefer video formats (vcodec != 'none')
            if f.get("vcodec") == "none":
                continue
            candidates.append(f)

        def score_format(fdict):
            # Score by (is_mp4, has_audio, height, tbr)
            is_mp4 = 1 if fdict.get("ext") == "mp4" else 0
            has_audio = 0 if fdict.get("acodec") == "none" else 1
            height = fdict.get("height") or 0
            tbr = fdict.get("tbr") or 0
            return (is_mp4, has_audio, height, tbr)

        if candidates:
            candidates.sort(key=score_format, reverse=True)
            stream_url = candidates[0].get("url")
        else:
            # no suitable format found
            available = []
            for f in formats:
                available.append({
                    "ext": f.get("ext"),
                    "protocol": f.get("protocol"),
                    "vcodec": f.get("vcodec"),
                    "acodec": f.get("acodec"),
                })
            return None, (
                "Could not resolve a playable stream URL from YouTube. "
                "No suitable progressive/video formats were found. Available formats: "
                f"{available}.\nConsider enabling an ffmpeg-based fallback if m3u8/HLS formats are returned."
            )

    # 2. Initialize OpenCV capture from stream URL
    def try_open(url):
        cap_local = cv2.VideoCapture(url)
        ok = cap_local.isOpened()
        if not ok:
            try:
                cap_local.release()
            except Exception:
                pass
            return None
        return cap_local

    cap = try_open(stream_url)
    # If opening failed and we have candidate formats, try them in order
    if cap is None and 'candidates' in locals() and candidates:
        for f in candidates[1:6]:
            url_try = f.get("url")
            cap = try_open(url_try)
            if cap is not None:
                stream_url = url_try
                break

    if cap is None:
        return None, (
            "Could not open video stream via OpenCV. The selected stream may be HLS/m3u8 or otherwise unsupported by your OpenCV/ffmpeg build. "
            "You can enable an ffmpeg fallback or run this locally with an OpenCV build that supports HLS."
        )

    # Setup Loop
    fps = cap.get(cv2.CAP_PROP_FPS) or 0.0
    # If FPS is unavailable or zero, sample by skipping a fixed number of frames
    skip_frames = int(round(fps * skip_seconds)) if fps > 0 else max(1, int(skip_seconds * 30))

    prev_frame = None
    slides = []

    # UI Feedback
    progress_bar = st.progress(0)
    status_text = st.empty()

    frames_processed = 0

    # If available, compute an estimated total frames for progress reporting
    total_frames_estimate = None
    try:
        if duration and fps > 0:
            total_frames_estimate = int(duration * fps)
    except Exception:
        total_frames_estimate = None

    while True:
        # Skip ahead to sample periodic frames
        for _ in range(skip_frames):
            cap.grab()

        success, frame = cap.read()
        if not success:
            break

        frames_processed += (skip_frames + 1)

        # Detection Logic: compare to previous sampled frame
        gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)

        is_new_slide = False
        if prev_frame is None:
            is_new_slide = True
        else:
            h, w = gray.shape
            diff = cv2.absdiff(prev_frame, gray)
            mse = float(np.sum(diff.astype("float64") ** 2) / (h * w))
            is_new_slide = mse > mse_threshold

        # Always update prev_frame to be the latest sampled frame
        prev_frame = gray

        if is_new_slide:
            rgb_frame = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
            slides.append(rgb_frame)
            status_text.text(f"Found {len(slides)} slides...")

        # Update progress bar using estimate if available
        if total_frames_estimate:
            current_progress = min(frames_processed / float(total_frames_estimate), 0.99)
        else:
            # fallback heuristic
            current_progress = min(frames_processed / 5000.0, 0.9)

        progress_bar.progress(current_progress)

        if frames_processed > max_frames:
            break

    cap.release()
    progress_bar.progress(1.0)
    status_text.text(f"Processing complete. Found {len(slides)} slides.")

    if not slides:
        return None, "No slides detected."

    # 3. Build PowerPoint
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_img in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Encode image to memory buffer (ensure bytes)
        is_success, buffer = cv2.imencode(".jpg", cv2.cvtColor(slide_img, cv2.COLOR_RGB2BGR))
        if not is_success:
            continue

        image_stream = io.BytesIO(buffer.tobytes())
        image_stream.seek(0)

        slide.shapes.add_picture(image_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)

    pptx_out = io.BytesIO()
    prs.save(pptx_out)
    pptx_out.seek(0)

    return pptx_out, None